from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from typing import Any

from django.db.models import QuerySet

from ..models import Campaign, CampaignProposal, ContentSubmission, InfluencerProfile


DEFAULT_EMV_RULES = {
    "base_cpm_eur": 18.0,
    "engagement_value_eur": 0.35,
    "format_multipliers": {
        "video": 1.25,
        "story": 0.85,
        "image": 1.0,
        "default": 1.0,
    },
}


def _to_float(value: Any, default: float = 0.0) -> float:
    try:
        if value is None:
            return default
        if isinstance(value, str):
            cleaned = value.replace("%", "").replace(" ", "").strip()
            if not cleaned:
                return default
            return float(cleaned)
        return float(value)
    except (TypeError, ValueError):
        return default


def _pick_stat(stats: dict[str, Any], keys: list[str], default: float = 0.0) -> float:
    for key in keys:
        if key in stats:
            value = _to_float(stats.get(key), default=default)
            if value > 0:
                return value
    return default


def _avg_social(profile: InfluencerProfile) -> tuple[float, float, list[str]]:
    networks = list(profile.social_networks.all())
    if not networks:
        return 0.0, 0.0, []
    followers = sum(float(n.followers_count or 0) for n in networks) / len(networks)
    engagement = sum(float(n.engagement_rate or 0) for n in networks) / len(networks)
    platforms = [str(n.platform or "").strip().lower() for n in networks if n.platform]
    return followers, engagement, platforms


def _jaccard(a: set[str], b: set[str]) -> float:
    if not a and not b:
        return 0.0
    inter = len(a & b)
    union = len(a | b)
    if union == 0:
        return 0.0
    return inter / union


def _proximity(a: float, b: float) -> float:
    max_v = max(abs(a), abs(b), 1.0)
    return max(0.0, 1.0 - (abs(a - b) / max_v))


@dataclass
class EmvRow:
    submission_id: int
    proposal_id: int
    influencer_id: int
    influencer_name: str
    impressions: float
    engagement: float
    emv_eur: float
    confidence: str


def classify_content_format(submission: ContentSubmission) -> str:
    url = str(submission.publication_url or "").lower()
    campaign_format = str(getattr(submission.proposal.campaign, "content_format", "") or "").lower()
    if any(token in url for token in ["reel", "video", "tiktok", "youtube", "shorts"]) or "video" in campaign_format:
        return "video"
    if "story" in url or "story" in campaign_format:
        return "story"
    return "image"


def compute_campaign_emv(campaign: Campaign) -> dict[str, Any]:
    submissions = ContentSubmission.objects.select_related(
        "proposal", "proposal__influencer", "proposal__influencer__user", "proposal__campaign"
    ).filter(proposal__campaign=campaign)

    rows: list[EmvRow] = []
    influencers_map: dict[int, dict[str, Any]] = {}

    for submission in submissions:
        stats = submission.final_stats or submission.initial_stats or {}
        if not isinstance(stats, dict):
            stats = {}

        impressions = _pick_stat(stats, ["impressions", "reach", "views", "view_count"], default=0.0)

        likes = _pick_stat(stats, ["likes", "like_count"], default=0.0)
        comments = _pick_stat(stats, ["comments", "comment_count"], default=0.0)
        shares = _pick_stat(stats, ["shares", "share_count"], default=0.0)
        saves = _pick_stat(stats, ["saves", "save_count"], default=0.0)
        clicks = _pick_stat(stats, ["clicks", "link_clicks"], default=0.0)
        engagement = likes + comments + shares + saves + clicks

        confidence = "high"

        if impressions <= 0:
            confidence = "medium"
            avg_views, avg_engagement, _ = _avg_social(submission.proposal.influencer)
            if avg_views > 0:
                impressions = avg_views
            elif avg_engagement > 0:
                impressions = max(500.0, avg_engagement * 120.0)
            else:
                confidence = "low"
                impressions = 1000.0

        if engagement <= 0:
            confidence = "medium" if confidence == "high" else confidence
            avg_followers, avg_engagement, _ = _avg_social(submission.proposal.influencer)
            inferred = impressions * max(avg_engagement, 0) / 100.0
            if inferred > 0:
                engagement = inferred
            elif avg_followers > 0:
                engagement = avg_followers * 0.01
            else:
                confidence = "low"
                engagement = max(10.0, impressions * 0.005)

        format_key = classify_content_format(submission)
        mult = DEFAULT_EMV_RULES["format_multipliers"].get(format_key, DEFAULT_EMV_RULES["format_multipliers"]["default"])

        emv = ((impressions / 1000.0) * DEFAULT_EMV_RULES["base_cpm_eur"] * mult) + (
            engagement * DEFAULT_EMV_RULES["engagement_value_eur"]
        )

        influencer = submission.proposal.influencer
        influencer_name = influencer.display_name or influencer.user.username

        row = EmvRow(
            submission_id=submission.id,
            proposal_id=submission.proposal_id,
            influencer_id=influencer.id,
            influencer_name=influencer_name,
            impressions=round(impressions, 2),
            engagement=round(engagement, 2),
            emv_eur=round(emv, 2),
            confidence=confidence,
        )
        rows.append(row)

        bucket = influencers_map.setdefault(
            influencer.id,
            {
                "influencer_id": influencer.id,
                "influencer_name": influencer_name,
                "submissions": 0,
                "impressions": 0.0,
                "engagement": 0.0,
                "emv_eur": 0.0,
            },
        )
        bucket["submissions"] += 1
        bucket["impressions"] += row.impressions
        bucket["engagement"] += row.engagement
        bucket["emv_eur"] += row.emv_eur

    budget_spent = 0.0
    paid = CampaignProposal.objects.filter(campaign=campaign, status="paid")
    for p in paid:
        budget_spent += float(p.escrow_amount or 0)

    total_emv = round(sum(r.emv_eur for r in rows), 2)
    ratio = round(total_emv / budget_spent, 3) if budget_spent > 0 else None

    influencers = sorted(influencers_map.values(), key=lambda x: x["emv_eur"], reverse=True)
    for item in influencers:
        item["impressions"] = round(item["impressions"], 2)
        item["engagement"] = round(item["engagement"], 2)
        item["emv_eur"] = round(item["emv_eur"], 2)

    confidence_score = 0.0
    if rows:
        confidence_map = {"high": 1.0, "medium": 0.65, "low": 0.35}
        confidence_score = round(sum(confidence_map.get(r.confidence, 0.35) for r in rows) / len(rows), 3)

    return {
        "campaign_id": campaign.id,
        "campaign_title": campaign.title,
        "currency": "EUR",
        "rules": DEFAULT_EMV_RULES,
        "submissions_count": len(rows),
        "emv_total_eur": total_emv,
        "budget_spent_eur": round(budget_spent, 2),
        "emv_vs_spend_ratio": ratio,
        "confidence_score": confidence_score,
        "by_influencer": influencers,
        "by_submission": [r.__dict__ for r in rows],
    }


def compute_lookalikes(
    campaign: Campaign,
    reference_influencer_id: int,
    limit: int = 20,
    min_score: float = 0.35,
) -> list[dict[str, Any]]:
    try:
        reference = InfluencerProfile.objects.select_related("user").prefetch_related("social_networks").get(pk=reference_influencer_id)
    except InfluencerProfile.DoesNotExist:
        return []

    ref_followers, ref_engagement, ref_platforms = _avg_social(reference)
    ref_themes = {str(v).strip().lower() for v in (reference.content_themes or []) if str(v).strip()}
    ref_location = str(reference.user.location or "").strip().lower()
    ref_rating = float(reference.average_rating or 0)

    qs: QuerySet[InfluencerProfile] = InfluencerProfile.objects.select_related("user").prefetch_related("social_networks")
    qs = qs.filter(onboarding_completed=True).exclude(pk=reference.id)

    platforms = [str(v).strip().lower() for v in (campaign.target_networks or []) if str(v).strip()]
    if platforms:
        qs = qs.filter(social_networks__platform__in=platforms).distinct()

    results: list[dict[str, Any]] = []

    for candidate in qs[:300]:
        c_followers, c_engagement, c_platforms = _avg_social(candidate)
        c_themes = {str(v).strip().lower() for v in (candidate.content_themes or []) if str(v).strip()}
        c_location = str(candidate.user.location or "").strip().lower()
        c_rating = float(candidate.average_rating or 0)

        theme_score = _jaccard(ref_themes, c_themes)
        platform_score = _jaccard(set(ref_platforms), set(c_platforms))
        followers_score = _proximity(ref_followers, c_followers)
        engagement_score = _proximity(ref_engagement, c_engagement)
        rating_score = _proximity(ref_rating, c_rating)
        location_score = 1.0 if (ref_location and c_location and ref_location == c_location) else 0.0

        total = (
            (0.30 * theme_score)
            + (0.20 * platform_score)
            + (0.20 * followers_score)
            + (0.15 * engagement_score)
            + (0.10 * rating_score)
            + (0.05 * location_score)
        )

        if total < min_score:
            continue

        reason_pairs = [
            ("themes", theme_score),
            ("platforms", platform_score),
            ("audience", followers_score),
            ("engagement", engagement_score),
            ("rating", rating_score),
            ("location", location_score),
        ]
        reason_pairs.sort(key=lambda x: x[1], reverse=True)

        results.append({
            "influencer_id": candidate.id,
            "display_name": candidate.display_name or candidate.user.username,
            "pseudo": candidate.display_name or candidate.user.username,
            "avatar": candidate.user.avatar.url if getattr(candidate.user, "avatar", None) else None,
            "location": candidate.user.location,
            "themes": list(candidate.content_themes or []),
            "platforms": c_platforms,
            "followers_avg": round(c_followers, 2),
            "engagement_rate_avg": round(c_engagement, 2),
            "rating": round(c_rating, 2),
            "score": round(total, 4),
            "reasons": [k for k, _ in reason_pairs[:3]],
        })

    results.sort(key=lambda x: x["score"], reverse=True)
    return results[: max(1, min(limit, 50))]


def build_campaign_report_payload(campaign: Campaign) -> dict[str, Any]:
    emv = compute_campaign_emv(campaign)
    proposals = CampaignProposal.objects.filter(campaign=campaign)

    status_counts: dict[str, int] = {}
    for status in [
        "pending",
        "accepted",
        "declined",
        "counter_offer",
        "contract_signed",
        "in_progress",
        "content_submitted",
        "validated",
        "paid",
        "disputed",
    ]:
        status_counts[status] = proposals.filter(status=status).count()

    return {
        "campaign": {
            "id": campaign.id,
            "title": campaign.title,
            "status": campaign.status,
            "deadline": campaign.deadline.isoformat() if campaign.deadline else None,
            "budget_per_influencer": float(campaign.price_per_influencer or 0),
            "max_influencers": int(campaign.max_influencers or 0),
        },
        "proposals": {
            "total": proposals.count(),
            "status_counts": status_counts,
        },
        "emv": emv,
    }


def render_report_pdf(report: dict[str, Any]) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas

    def fmt_currency(value: Any) -> str:
        return f"EUR {float(value or 0):,.2f}".replace(",", " ")

    def fmt_ratio(value: Any) -> str:
        if value is None:
            return "-"
        return f"{float(value):.2f}x"

    def status_label(raw: str) -> str:
        return str(raw or "-").replace("_", " ").title()

    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4

    campaign = report["campaign"]
    emv = report["emv"]
    proposals = report.get("proposals", {})

    # Background inspired by transactional email template.
    c.setFillColor(colors.HexColor("#EEF2FF"))
    c.rect(0, 0, width, height, stroke=0, fill=1)

    card_x, card_y = 24, 24
    card_w, card_h = width - 48, height - 48
    c.setFillColor(colors.white)
    c.roundRect(card_x, card_y, card_w, card_h, 16, stroke=0, fill=1)

    # Header band
    header_h = 94
    c.setFillColor(colors.HexColor("#2563EB"))
    c.roundRect(card_x, card_y + card_h - header_h, card_w, header_h, 16, stroke=0, fill=1)
    # Square-off lower corners visually.
    c.rect(card_x, card_y + card_h - header_h, card_w, 20, stroke=0, fill=1)

    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 9)
    c.drawString(card_x + 20, card_y + card_h - 22, "INFLUCONNECT")
    c.setFont("Helvetica-Bold", 18)
    c.drawString(card_x + 20, card_y + card_h - 44, "Campaign Performance Report")
    c.setFont("Helvetica", 11)
    c.drawString(card_x + 20, card_y + card_h - 62, f"{campaign['title']}  |  #{campaign['id']}")

    y = card_y + card_h - header_h - 24

    # Meta line.
    c.setFillColor(colors.HexColor("#334155"))
    c.setFont("Helvetica", 10)
    c.drawString(card_x + 20, y, f"Status: {status_label(campaign.get('status'))}")
    c.drawString(card_x + 190, y, f"Deadline: {campaign.get('deadline') or '-'}")
    c.drawString(card_x + 360, y, f"Budget / influencer: {fmt_currency(campaign.get('budget_per_influencer'))}")
    y -= 30

    # KPI cards.
    kpi_labels = [
        ("Total EMV", fmt_currency(emv.get("emv_total_eur"))),
        ("Budget spent", fmt_currency(emv.get("budget_spent_eur"))),
        ("EMV / Spend", fmt_ratio(emv.get("emv_vs_spend_ratio"))),
        ("Confidence", f"{float(emv.get('confidence_score') or 0):.2f}"),
    ]
    kpi_w = (card_w - 40 - 30) / 4
    kpi_h = 62
    for idx, (label, value) in enumerate(kpi_labels):
        x = card_x + 20 + idx * (kpi_w + 10)
        c.setFillColor(colors.HexColor("#F8FAFC"))
        c.roundRect(x, y - kpi_h, kpi_w, kpi_h, 9, stroke=0, fill=1)
        c.setFillColor(colors.HexColor("#64748B"))
        c.setFont("Helvetica", 9)
        c.drawString(x + 10, y - 18, label)
        c.setFillColor(colors.HexColor("#0F172A"))
        c.setFont("Helvetica-Bold", 13)
        c.drawString(x + 10, y - 38, value)
    y -= kpi_h + 24

    # Status distribution
    c.setFillColor(colors.HexColor("#0F172A"))
    c.setFont("Helvetica-Bold", 12)
    c.drawString(card_x + 20, y, "Proposal pipeline")
    y -= 16
    c.setFont("Helvetica", 10)
    status_counts = proposals.get("status_counts", {}) or {}
    ordered_status = [
        "pending", "counter_offer", "accepted", "contract_signed", "in_progress",
        "content_submitted", "validated", "paid", "declined", "disputed",
    ]
    left_x = card_x + 20
    right_x = card_x + card_w / 2 + 10
    row_h = 14
    row_idx = 0
    for key in ordered_status:
        count = int(status_counts.get(key, 0) or 0)
        target_x = left_x if row_idx < 5 else right_x
        target_y = y - (row_idx % 5) * row_h
        c.setFillColor(colors.HexColor("#334155"))
        c.drawString(target_x, target_y, f"{status_label(key)}")
        c.setFillColor(colors.HexColor("#0F172A"))
        c.setFont("Helvetica-Bold", 10)
        c.drawRightString(target_x + 170, target_y, str(count))
        c.setFont("Helvetica", 10)
        row_idx += 1

    y -= 84

    # Top influencers table.
    c.setFillColor(colors.HexColor("#0F172A"))
    c.setFont("Helvetica-Bold", 12)
    c.drawString(card_x + 20, y, "Top influencers by EMV")
    y -= 16

    rows = emv.get("by_influencer", [])[:8]
    table_x = card_x + 20
    table_w = card_w - 40
    header_y = y
    c.setFillColor(colors.HexColor("#E2E8F0"))
    c.roundRect(table_x, header_y - 16, table_w, 16, 4, stroke=0, fill=1)
    c.setFillColor(colors.HexColor("#0F172A"))
    c.setFont("Helvetica-Bold", 9)
    c.drawString(table_x + 8, header_y - 11, "Influencer")
    c.drawString(table_x + 250, header_y - 11, "Submissions")
    c.drawString(table_x + 340, header_y - 11, "Impressions")
    c.drawString(table_x + 450, header_y - 11, "EMV")

    row_y = header_y - 24
    c.setFont("Helvetica", 9)
    for r in rows:
        if row_y < card_y + 24:
            c.showPage()
            c.setFillColor(colors.HexColor("#EEF2FF"))
            c.rect(0, 0, width, height, stroke=0, fill=1)
            c.setFillColor(colors.white)
            c.roundRect(card_x, card_y, card_w, card_h, 16, stroke=0, fill=1)
            row_y = card_y + card_h - 40
        c.setFillColor(colors.HexColor("#334155"))
        c.drawString(table_x + 8, row_y, str(r.get("influencer_name") or "-"))
        c.drawRightString(table_x + 322, row_y, str(int(r.get("submissions") or 0)))
        c.drawRightString(table_x + 432, row_y, f"{float(r.get('impressions') or 0):,.0f}".replace(",", " "))
        c.drawRightString(table_x + 540, row_y, fmt_currency(r.get("emv_eur")))
        c.setStrokeColor(colors.HexColor("#E2E8F0"))
        c.line(table_x, row_y - 4, table_x + table_w, row_y - 4)
        row_y -= 16

    c.setFillColor(colors.HexColor("#64748B"))
    c.setFont("Helvetica", 8)
    c.drawString(card_x + 20, card_y + 12, "Generated by InfluConnect reporting")

    c.save()
    buffer.seek(0)
    return buffer.read()


def render_report_pptx(report: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    def fmt_currency(value: Any) -> str:
        return f"EUR {float(value or 0):,.2f}".replace(",", " ")

    def status_label(raw: str) -> str:
        return str(raw or "-").replace("_", " ").title()

    def add_bg(slide, color_hex: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex)

    def add_header(slide, title: str, subtitle: str):
        bar = slide.shapes.add_shape(1, Inches(0.3), Inches(0.2), Inches(12.7), Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string("2563EB")
        bar.line.fill.background()
        tf = bar.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor.from_string("FFFFFF")
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor.from_string("DBEAFE")

    def add_kpi_box(slide, x: float, y: float, title: str, value: str):
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.95), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor.from_string("F8FAFC")
        box.line.color.rgb = RGBColor.from_string("E2E8F0")
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor.from_string("64748B")
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.bold = True
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor.from_string("0F172A")

    prs = Presentation()
    campaign = report["campaign"]
    proposals = report["proposals"]
    emv = report["emv"]

    # Slide 1 - cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "EEF2FF")
    slide.shapes.title.text = "Campaign Performance Report"
    slide.placeholders[1].text = f"{campaign['title']}  |  #{campaign['id']}"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("1E3A8A")
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)

    # Slide 2 - KPI overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bg(slide, "EEF2FF")
    add_header(slide, "Campaign KPI overview", f"Status: {status_label(campaign['status'])}   Deadline: {campaign['deadline'] or '-'}")
    slide.shapes.placeholders[1].text_frame.clear()
    add_kpi_box(slide, 0.6, 1.6, "Total proposals", str(proposals["total"]))
    add_kpi_box(slide, 3.8, 1.6, "Budget / influencer", fmt_currency(campaign.get("budget_per_influencer")))
    add_kpi_box(slide, 7.0, 1.6, "Max influencers", str(campaign.get("max_influencers") or 0))
    add_kpi_box(slide, 10.2, 1.6, "Campaign ID", f"#{campaign['id']}")

    status_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(12.0), Inches(2.9))
    tf = status_box.text_frame
    tf.clear()
    title_p = tf.paragraphs[0]
    title_p.text = "Pipeline distribution"
    title_p.font.bold = True
    title_p.font.size = Pt(16)
    title_p.font.color.rgb = RGBColor.from_string("0F172A")

    status_counts = proposals.get("status_counts", {}) or {}
    ordered_status = [
        "pending", "counter_offer", "accepted", "contract_signed", "in_progress",
        "content_submitted", "validated", "paid", "declined", "disputed",
    ]
    for key in ordered_status:
        p = tf.add_paragraph()
        p.text = f"{status_label(key)}: {int(status_counts.get(key, 0) or 0)}"
        p.font.size = Pt(12)
        p.level = 1

    # Slide 3 - EMV summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bg(slide, "EEF2FF")
    add_header(slide, "Earned Media Value", "Estimated value and campaign efficiency")
    slide.shapes.placeholders[1].text_frame.clear()
    add_kpi_box(slide, 0.6, 1.6, "Total EMV", fmt_currency(emv["emv_total_eur"]))
    add_kpi_box(slide, 3.8, 1.6, "Budget spent", fmt_currency(emv["budget_spent_eur"]))
    ratio = emv.get("emv_vs_spend_ratio")
    add_kpi_box(slide, 7.0, 1.6, "EMV / Spend", (f"{float(ratio):.2f}x" if ratio is not None else "-"))
    add_kpi_box(slide, 10.2, 1.6, "Confidence", f"{float(emv.get('confidence_score', 0)):.2f}")

    note = slide.shapes.add_textbox(Inches(0.8), Inches(3.55), Inches(12.0), Inches(2.6))
    nt = note.text_frame
    nt.clear()
    p = nt.paragraphs[0]
    p.text = "Method"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor.from_string("0F172A")
    p2 = nt.add_paragraph()
    p2.text = "EMV combines impressions (CPM-based) and engagement interactions with format multipliers."
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor.from_string("334155")
    p3 = nt.add_paragraph()
    p3.text = "Confidence decreases when performance metrics are inferred from fallback social averages."
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor.from_string("334155")

    # Slide 4 - top influencers
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bg(slide, "EEF2FF")
    slide.shapes.title.text = "Top influencers by EMV"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("1E293B")
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)
    rows = emv.get("by_influencer", [])[:8]

    if rows:
        table = slide.shapes.add_table(len(rows) + 1, 4, left=Inches(0.7), top=Inches(1.5), width=Inches(12.0), height=Inches(4.7)).table
        table.cell(0, 0).text = "Influencer"
        table.cell(0, 1).text = "Submissions"
        table.cell(0, 2).text = "Impressions"
        table.cell(0, 3).text = "EMV (EUR)"

        for col in range(4):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("DBEAFE")
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor.from_string("1E3A8A")
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        for idx, row in enumerate(rows, start=1):
            table.cell(idx, 0).text = str(row.get("influencer_name", "-"))
            table.cell(idx, 1).text = str(row.get("submissions", 0))
            table.cell(idx, 2).text = f"{float(row.get('impressions', 0)):.0f}"
            table.cell(idx, 3).text = f"{float(row.get('emv_eur', 0)):.2f}"
            for col in range(4):
                run = table.cell(idx, col).text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor.from_string("0F172A")

    footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(8.0), Inches(0.4))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Generated by InfluConnect reporting"
    fp.font.size = Pt(10)
    fp.font.color.rgb = RGBColor.from_string("64748B")

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
